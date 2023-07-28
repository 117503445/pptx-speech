import collections.abc  # for high Python3 version
from pptx import Presentation

from pathlib import Path


def get_notes(file_pptx: Path) -> list[str]:
    ppt = Presentation(str(file_pptx.absolute()))

    notes = []

    for slide in ppt.slides:
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append(textNote)

    return handle_animation(notes)

def handle_animation(notes: list[str]) -> list[str]:
    # input
    # ["page1", "page2", "$speech$ page3-with-animation $click$ click-once $click$ click-twice", "$speech$ page3-with-animation $click$ click-once $click$ click-twice", "$speech$ page3-with-animation $click$ click-once $click$ click-twice", "$speech$ page3-with-animation $click$ click-once $click$ click-twice", "page end"]

    # output
    # ["page1", "page2", "page3-with-animation", "click-once", "click-twice", "page end"]

    notes_new = []

    if not notes:
        return notes_new
    
    i = 0
    while i < len(notes):
        note = notes[i]

        if '$speech$' in note:
            note = note.replace("$speech$", "")
            note = note.split("$click$")
            notes_new.extend(note)
            i += len(note) + 1
        else:
            notes_new.append(note)
            i += 1

    return notes_new