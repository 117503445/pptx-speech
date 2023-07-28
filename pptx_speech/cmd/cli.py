import collections.abc # for high Python3 version
from pptx import Presentation

from htutil import file

f = 'example.pptx'

ppt=Presentation(f)

notes = []

for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note. 
    textNote = slide.notes_slide.notes_text_frame.text
    notes.append((page,textNote)) 

print(notes)
file.write_json('1.json', notes)